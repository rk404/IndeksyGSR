"""
Generuje prezentację PPTX dla IndeksyGSR.
Odbiorcy: studenci wydziału okrętownictwa (nieIT).

Uruchomienie:
    uv run python scripts/generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Kolory Solarized Dark ──────────────────────────────────────────────────
BG     = RGBColor(0x00, 0x2b, 0x36)
PANEL  = RGBColor(0x07, 0x36, 0x42)
PANEL2 = RGBColor(0x0a, 0x42, 0x4f)
CYAN   = RGBColor(0x2a, 0xa1, 0x98)
BLUE   = RGBColor(0x26, 0x8b, 0xd2)
TEXT   = RGBColor(0xfd, 0xf6, 0xe3)
MUTED  = RGBColor(0x93, 0xa1, 0xa1)
GREEN  = RGBColor(0x85, 0x99, 0x00)
ORANGE = RGBColor(0xcb, 0x4b, 0x16)
YELLOW = RGBColor(0xb5, 0x89, 0x00)

# ── Wymiary slajdu ─────────────────────────────────────────────────────────
SW = Inches(13.33)   # szerokość slajdu
SH = Inches(7.5)     # wysokość slajdu

# ── Stałe marginesów — stosowane we WSZYSTKICH slajdach ───────────────────
ML  = Inches(0.55)              # lewy margines
MR  = Inches(0.55)              # prawy margines
CW  = SW - ML - MR              # szerokość obszaru treści = 12.23"
MT  = Inches(1.35)              # górna granica treści (pod nagłówkiem + linią)
MB  = Inches(0.35)              # dolny margines
BOT = SH - MB                   # dolna granica treści = 7.15"
CH  = BOT - MT                  # wysokość obszaru treści = 5.8"

# Wysokości standardowych elementów
HEAD_H   = Inches(0.65)         # wysokość nagłówka
HEAD_Y   = Inches(0.28)         # y nagłówka
DIV_Y    = Inches(1.05)         # y linii poziomej


# ── Narzędzia pomocnicze ───────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color and border_pt > 0:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def heading(slide, text):
    """Turkusowy nagłówek — zawsze na tej samej pozycji."""
    add_text(slide, text, ML, HEAD_Y, CW, HEAD_H,
             font_size=34, bold=True, color=CYAN, align=PP_ALIGN.LEFT)


def divider(slide):
    """Turkusowa linia pozioma — zawsze na tej samej pozycji."""
    add_rect(slide, ML, DIV_Y, CW, Pt(2), CYAN)


def flow_box(slide, x, y, w, h, text, fill=PANEL, text_color=TEXT,
             font_size=14, border_color=CYAN):
    add_rect(slide, x, y, w, h, fill, border_color, 1.5)
    add_text(slide, text, x, y, w, h,
             font_size=font_size, color=text_color,
             align=PP_ALIGN.CENTER, wrap=True)


def arrow(slide, x1, y1, x2, y2):
    c = slide.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = CYAN
    c.line.width = Pt(2)


def stat_card(slide, x, y, w, h, number, label, color=CYAN):
    add_rect(slide, x, y, w, h, PANEL, color, 1.5)
    add_text(slide, number, x, y + Inches(0.15), w, Inches(0.8),
             font_size=42, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + Inches(1.0), w, Inches(0.55),
             font_size=13, color=MUTED, align=PP_ALIGN.CENTER)


def module_tile(slide, x, y, icon, title, desc, color=BLUE):
    TW = Inches(5.95)
    TH = Inches(1.6)
    add_rect(slide, x, y, TW, TH, PANEL, color, 1.5)
    add_text(slide, f"{icon}  {title}",
             x + Inches(0.2), y + Inches(0.1), TW - Inches(0.3), Inches(0.5),
             font_size=17, bold=True, color=color)
    add_text(slide, desc,
             x + Inches(0.2), y + Inches(0.62), TW - Inches(0.3), Inches(0.88),
             font_size=13, color=MUTED)


# ── SLAJDY ────────────────────────────────────────────────────────────────

def slide_01_title(prs):
    s = blank_slide(prs)
    bg(s)
    # Pasek pionowy
    add_rect(s, Inches(0), Inches(0), Inches(0.22), SH, CYAN)
    # Tytuł
    add_text(s, "IndeksyGSR",
             ML, Inches(1.7), Inches(10), Inches(1.35),
             font_size=72, bold=True, color=CYAN)
    # Podtytuł
    add_text(s, "Inteligentny system wyszukiwania indeksów materiałowych",
             ML, Inches(3.15), CW, Inches(0.75),
             font_size=26, color=TEXT)
    # Linia
    add_rect(s, ML, Inches(4.05), Inches(8), Pt(2), BLUE)
    # Info
    add_text(s, "Projekt dyplomowy · Politechnika Gdańska",
             ML, Inches(4.3), CW, Inches(0.5),
             font_size=16, color=MUTED)
    # Rok
    add_text(s, "2025 / 2026",
             SW - Inches(2.8), BOT - Inches(0.45), Inches(2.5), Inches(0.4),
             font_size=14, color=MUTED, align=PP_ALIGN.RIGHT)


def slide_02_context(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Czym są indeksy materiałów?")
    divider(s)

    add_text(s,
             "Indeks materiału to unikalny kod identyfikujący każdy rodzaj materiału lub części "
             "w magazynie lub katalogu technicznym — podobnie jak ISBN dla książki.",
             ML, MT, CW, Inches(0.8),
             font_size=16, color=TEXT)

    # 3 karty indeksów — równo podzielone w obszarze CW
    GAP   = Inches(0.2)
    TW    = (CW - 2 * GAP) / 3
    TH    = Inches(2.1)
    TY    = MT + Inches(1.0)
    cards = [
        ("ŚRUBA-M20-DIN931-OGNIOWO", "Śruba M20×80 DIN 931\nocynkowana ogniowo", CYAN),
        ("RURA-DN50-ST35-BW",        "Rura stalowa DN50\nbezszwowa St35",         BLUE),
        ("USZCZ-GUMA-EPDM-3MM",      "Uszczelka gumowa EPDM\ngrubość 3 mm",       GREEN),
    ]
    for i, (code, desc, color) in enumerate(cards):
        x = ML + i * (TW + GAP)
        add_rect(s, x, TY, TW, TH, PANEL, color, 2)
        add_text(s, code, x + Inches(0.15), TY + Inches(0.1), TW - Inches(0.2), Inches(0.5),
                 font_size=11, bold=True, color=color)
        add_rect(s, x, TY + Inches(0.65), TW, Pt(1), color)
        add_text(s, desc, x + Inches(0.15), TY + Inches(0.75), TW - Inches(0.2), Inches(1.2),
                 font_size=14, color=TEXT)

    # Dolny fakt
    FY = TY + TH + Inches(0.2)
    FH = BOT - FY
    add_rect(s, ML, FY, CW, FH, PANEL2, YELLOW, 1.5)
    add_text(s,
             "🏭  W dużych stoczniach i przedsiębiorstwach przemysłowych katalogi materiałów "
             "liczą dziesiątki tysięcy pozycji.\n"
             "Baza IndeksyGSR zawiera  69 000  indeksów.",
             ML + Inches(0.25), FY + Inches(0.15), CW - Inches(0.4), FH - Inches(0.2),
             font_size=15, color=TEXT)


def slide_03_problem(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Problem — jak wygląda to dziś?")
    divider(s)

    # 3 boxy przepływu — równo w CW
    GAP   = Inches(0.2)
    BW    = (CW - 2 * GAP) / 3
    BH    = Inches(2.3)
    BY    = MT
    boxes = [
        ("📋 Pracownik otrzymuje\nopis produktu:\n\"śruba M20\nocynkowana ogniowo\"", BLUE),
        ("📚 Przeszukuje ręcznie\nkatalog\n69 000 pozycji",                           ORANGE),
        ("⏱ Wynik:\nCzasochłonne\nPodatne na błędy\nNiespójne",                      ORANGE),
    ]
    for i, (text, border) in enumerate(boxes):
        x = ML + i * (BW + GAP)
        flow_box(s, x, BY, BW, BH, text, PANEL, TEXT, 15, border)

    # Strzałki
    mid = BY + BH / 2
    for i in range(2):
        x_end = ML + (i + 1) * (BW + GAP)
        arrow(s, x_end - GAP, mid, x_end, mid)

    # Blok skutków
    EY = BY + BH + Inches(0.2)
    EH = BOT - EY
    add_rect(s, ML, EY, CW, EH, PANEL, ORANGE, 1)
    add_text(s, "Skutki dla organizacji:",
             ML + Inches(0.25), EY + Inches(0.12), CW - Inches(0.4), Inches(0.38),
             font_size=14, bold=True, color=ORANGE)
    problems = [
        "❌  Ten sam materiał może być wpisany pod różnymi kodami przez różnych pracowników",
        "❌  Ręczne przeszukiwanie katalogu 69 000 pozycji zajmuje kilkanaście minut",
        "❌  Synonimy i skróty (np. \"ocynkowana ogniowo\" = \"OGNIOWO\") nie są rozpoznawane",
        "❌  Błędy w przypisaniu indeksu powodują problemy z zamówieniami i magazynowaniem",
    ]
    line_h = (EH - Inches(0.55)) / len(problems)
    for i, p in enumerate(problems):
        add_text(s, p,
                 ML + Inches(0.25), EY + Inches(0.55) + i * line_h,
                 CW - Inches(0.4), line_h,
                 font_size=13, color=TEXT)


def slide_04_goal(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Cel projektu")
    divider(s)

    GAP = Inches(0.25)
    COL = (CW - GAP) / 2
    COL_H = BOT - MT

    # PRZED
    add_rect(s, ML, MT, COL, COL_H, PANEL, ORANGE, 2)
    add_text(s, "PRZED", ML, MT + Inches(0.1), COL, Inches(0.45),
             font_size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_rect(s, ML, MT + Inches(0.6), COL, Pt(1.5), ORANGE)
    before_items = [
        "Pracownik wpisuje opis",
        "Ręcznie przeszukuje Excel / katalog",
        "Kilkanaście minut pracy",
        "Ryzyko błędu lub pominięcia",
        "Brak pewności co do indeksu",
    ]
    lh = (COL_H - Inches(0.8)) / len(before_items)
    for i, item in enumerate(before_items):
        add_text(s, f"  {item}",
                 ML + Inches(0.2), MT + Inches(0.75) + i * lh,
                 COL - Inches(0.3), lh,
                 font_size=14, color=MUTED)

    # Strzałka
    add_text(s, "→",
             ML + COL + Inches(0.02), MT + COL_H / 2 - Inches(0.3),
             GAP, Inches(0.6),
             font_size=32, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

    # PO
    x2 = ML + COL + GAP
    add_rect(s, x2, MT, COL, COL_H, PANEL, GREEN, 2)
    add_text(s, "PO WDROŻENIU IndeksyGSR", x2, MT + Inches(0.1), COL, Inches(0.45),
             font_size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_rect(s, x2, MT + Inches(0.6), COL, Pt(1.5), GREEN)
    after_items = [
        "Pracownik wpisuje opis",
        "System zwraca wyniki w < 1 sekundy",
        "Wyniki posortowane wg dopasowania",
        "Wyświetlany % pewności dopasowania",
        "Możliwość zaproponowania nowego indeksu",
    ]
    for i, item in enumerate(after_items):
        add_text(s, f"✓  {item}",
                 x2 + Inches(0.2), MT + Inches(0.75) + i * lh,
                 COL - Inches(0.3), lh,
                 font_size=14, color=TEXT)


def slide_05_how(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Jak działa system?")
    divider(s)

    # Baner analogii
    AH = Inches(0.7)
    add_rect(s, ML, MT, CW, AH, PANEL2, BLUE, 1)
    add_text(s,
             "💡  Analogia: Google rozumie że \"duży metalowy statek\" = \"okręt\"."
             "  IndeksyGSR rozumie że \"ocynkowana ogniowo\" = \"OGNIOWO\", \"nierdzewna\" = \"A2\".",
             ML + Inches(0.2), MT + Inches(0.1), CW - Inches(0.3), AH - Inches(0.1),
             font_size=13, color=CYAN)

    # Diagram 4 kroków
    GAP   = Inches(0.15)
    FBW   = (CW - 3 * GAP) / 4
    FBH   = Inches(1.65)
    FBY   = MT + AH + Inches(0.2)
    steps = [
        ('"śruba M20\nocynkowana ogniowo"', PANEL, TEXT),
        ('Zrozumienie\nznaczenia\ni synonimów',  PANEL2, CYAN),
        ('Przeszukanie\nbazy 69 000\nindeksów',  PANEL2, BLUE),
        ('Wyniki\nz oceną\npewności',            PANEL, GREEN),
    ]
    xs = [ML + i * (FBW + GAP) for i in range(4)]
    for x, (text, fill, col) in zip(xs, steps):
        flow_box(s, x, FBY, FBW, FBH, text, fill, col, 13, col)

    mid = FBY + FBH / 2
    for i in range(3):
        arrow(s, xs[i] + FBW, mid, xs[i + 1], mid)

    # Przykładowe wyniki
    LY = FBY + FBH + Inches(0.2)
    add_text(s, "Przykładowe wyniki dla zapytania \"śruba M20 ocynkowana ogniowo\":",
             ML, LY, CW, Inches(0.35),
             font_size=13, bold=True, color=MUTED)

    results = [
        ("ŚRUBA-M20-DIN931-OGNIOWO",  "Śruba M20×80 ocynkowana ogniowo", "94%", GREEN),
        ("ŚRUBA-M20-DIN931-A2",       "Śruba M20×80 nierdzewna A2",      "71%", YELLOW),
        ("ŚRUBA-M20-DIN933-OGNIOWO",  "Śruba M20×50 ocynkowana ogniowo", "58%", ORANGE),
    ]
    RH = (BOT - LY - Inches(0.4)) / len(results)
    for i, (code, name, pct, color) in enumerate(results):
        ry = LY + Inches(0.4) + i * RH
        add_rect(s, ML, ry, CW, RH - Inches(0.04), PANEL)
        add_text(s, code,  ML + Inches(0.2),         ry + Inches(0.05), Inches(4.5), RH, font_size=12, bold=True, color=CYAN)
        add_text(s, name,  ML + Inches(5.0),          ry + Inches(0.05), Inches(5.5), RH, font_size=12, color=TEXT)
        add_text(s, pct,   ML + CW - Inches(1.2),    ry + Inches(0.05), Inches(1.0), RH, font_size=13, bold=True, color=color, align=PP_ALIGN.RIGHT)


def slide_06_modules(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Moduły aplikacji")
    divider(s)

    add_text(s, "Aplikacja posiada 4 tryby pracy dostępne przez przeglądarkę internetową:",
             ML, MT, CW, Inches(0.45),
             font_size=15, color=MUTED)

    GAP   = Inches(0.25)
    TW    = (CW - GAP) / 2
    TH    = Inches(1.65)
    tiles = [
        ("🔍", "Wyszukiwanie po tekście",
         "Wpisz opis materiału. System zwraca listę pasujących indeksów z oceną pewności dopasowania.",
         CYAN,   ML,        MT + Inches(0.55)),
        ("🌐", "Wyszukiwanie po URL",
         "Wklej link do produktu w sklepie. System automatycznie pobiera dane ze strony i szuka indeksu.",
         BLUE,   ML + TW + GAP, MT + Inches(0.55)),
        ("📋", "Wyszukiwanie zbiorcze",
         "Wgraj plik Excel z listą opisów. Pobierz gotowy plik z przypisanymi indeksami i % pewności.",
         GREEN,  ML,        MT + Inches(0.55) + TH + GAP),
        ("📝", "Propozycje nowych indeksów",
         "Jeśli żaden indeks nie pasuje, zaproponuj nowy. System podpowiada kategorię i zatwierdza do bazy.",
         YELLOW, ML + TW + GAP, MT + Inches(0.55) + TH + GAP),
    ]
    for icon, title, desc, color, x, y in tiles:
        add_rect(s, x, y, TW, TH, PANEL, color, 1.5)
        add_text(s, f"{icon}  {title}",
                 x + Inches(0.2), y + Inches(0.1), TW - Inches(0.3), Inches(0.5),
                 font_size=17, bold=True, color=color)
        add_text(s, desc,
                 x + Inches(0.2), y + Inches(0.65), TW - Inches(0.3), TH - Inches(0.7),
                 font_size=13, color=MUTED)


def slide_07_search_text(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Moduł 1 — Wyszukiwanie po tekście")
    divider(s)

    # Ramka okna przeglądarki
    WIN_H = BOT - MT
    add_rect(s, ML, MT, CW, WIN_H, PANEL, CYAN, 1.5)

    # Pasek adresu
    BAR_H = Inches(0.38)
    add_rect(s, ML, MT, CW, BAR_H, PANEL2)
    add_text(s, "🔒  indeksygsr.app/search",
             ML + Inches(0.2), MT + Inches(0.04), Inches(7), BAR_H - Inches(0.05),
             font_size=11, color=MUTED)

    # Pole wyszukiwania
    SY = MT + BAR_H + Inches(0.15)
    SH2 = Inches(0.52)
    add_rect(s, ML + Inches(0.25), SY, Inches(9.8), SH2, BG, CYAN, 1.5)
    add_text(s, "śruba M20 ocynkowana ogniowo 80mm DIN 931",
             ML + Inches(0.45), SY + Inches(0.07), Inches(9.4), SH2,
             font_size=13, color=TEXT)
    add_rect(s, ML + Inches(10.2), SY, Inches(1.75), SH2, BLUE)
    add_text(s, "Szukaj",
             ML + Inches(10.2), SY + Inches(0.07), Inches(1.75), SH2,
             font_size=13, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

    # Nagłówek tabeli
    TY = SY + SH2 + Inches(0.15)
    HDR_H = Inches(0.35)
    add_rect(s, ML + Inches(0.25), TY, CW - Inches(0.5), HDR_H, PANEL2)
    add_text(s, "Nazwa / Indeks",
             ML + Inches(0.85), TY + Inches(0.04), Inches(6.5), HDR_H,
             font_size=10, bold=True, color=MUTED)
    add_text(s, "Pewność dopasowania",
             ML + Inches(8.2), TY + Inches(0.04), Inches(3.5), HDR_H,
             font_size=10, bold=True, color=MUTED)

    # Wiersze
    results = [
        ("ŚRUBA-M20-DIN931-OGNIOWO",     "Śruba M20×80 DIN 931 ocynkowana ogniowo",  "94%", GREEN,  True),
        ("ŚRUBA-M20-DIN931-OGNIOWO-100", "Śruba M20×100 DIN 931 ocynkowana ogniowo", "87%", GREEN,  False),
        ("ŚRUBA-M20-DIN933-OGNIOWO",     "Śruba M20×50 DIN 933 ocynkowana ogniowo",  "72%", YELLOW, False),
        ("ŚRUBA-M20-IS04014-OGNIOWO",    "Śruba M20×60 ISO 4014 ocynkowana ogniowo", "65%", YELLOW, False),
    ]
    avail_h = BOT - (TY + HDR_H) - Inches(0.6)
    RH = avail_h / len(results)
    for i, (code, name, pct, color, checked) in enumerate(results):
        ry = TY + HDR_H + i * RH
        add_rect(s, ML + Inches(0.25), ry, CW - Inches(0.5), RH - Inches(0.02),
                 PANEL2 if checked else PANEL)
        chk = "☑" if checked else "☐"
        add_text(s, chk,  ML + Inches(0.35), ry + Inches(0.05), Inches(0.35), RH, font_size=13, color=CYAN if checked else MUTED)
        add_text(s, code, ML + Inches(0.8),  ry + Inches(0.05), Inches(3.8),  RH, font_size=11, bold=True, color=CYAN)
        add_text(s, name, ML + Inches(4.8),  ry + Inches(0.05), Inches(5.5),  RH, font_size=11, color=TEXT)
        add_text(s, pct,  ML + CW - Inches(1.3), ry + Inches(0.05), Inches(1.0), RH, font_size=12, bold=True, color=color, align=PP_ALIGN.RIGHT)

    # Przycisk Zapisz
    BY2 = BOT - Inches(0.52)
    add_rect(s, ML + Inches(0.25), BY2, Inches(2.9), Inches(0.46), GREEN)
    add_text(s, "✓  Zapisz zaznaczone (1)",
             ML + Inches(0.25), BY2 + Inches(0.04), Inches(2.9), Inches(0.38),
             font_size=12, bold=True, color=TEXT, align=PP_ALIGN.CENTER)


def slide_08_url(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Moduł 2 — Wyszukiwanie po URL sklepu")
    divider(s)

    add_text(s,
             "Użytkownik wkleja link do produktu w sklepie internetowym (np. Allegro, TME). "
             "System automatycznie pobiera dane ze strony — bez konieczności ręcznego przepisywania opisu.",
             ML, MT, CW, Inches(0.75),
             font_size=15, color=TEXT)

    # Diagram 4 kroków
    GAP  = Inches(0.2)
    FBW  = (CW - 3 * GAP) / 4
    FBH  = Inches(1.7)
    FBY  = MT + Inches(0.9)
    steps = [
        ("🌐 Link\ndo produktu\nw sklepie",      PANEL,  TEXT),
        ("🤖 Robot\npobiera dane\nze strony",    PANEL2, CYAN),
        ("📝 Ekstrakcja:\ntytuł, opis,\nspecyfikacja", PANEL2, BLUE),
        ("✅ Lista\npasujących\nindeksów",        PANEL,  GREEN),
    ]
    xs = [ML + i * (FBW + GAP) for i in range(4)]
    for x, (text, fill, col) in zip(xs, steps):
        flow_box(s, x, FBY, FBW, FBH, text, fill, col, 13, col)

    mid = FBY + FBH / 2
    for i in range(3):
        arrow(s, xs[i] + FBW, mid, xs[i + 1], mid)

    # Przykład URL
    UY = FBY + FBH + Inches(0.2)
    UH = Inches(0.52)
    add_rect(s, ML, UY, CW, UH, PANEL2, BLUE, 1)
    add_text(s, "Przykład: https://allegro.pl/oferta/sruba-m20-ocynkowana-ogniowo-din931-...",
             ML + Inches(0.25), UY + Inches(0.1), CW - Inches(0.4), UH - Inches(0.1),
             font_size=12, color=CYAN)

    # Obsługiwane sklepy
    add_text(s, "Obsługiwane źródła:",
             ML, UY + UH + Inches(0.2), CW, Inches(0.38),
             font_size=13, bold=True, color=MUTED)
    shops = ["🛒  Allegro.pl", "⚡  TME.eu", "🌍  Dowolna strona produktowa"]
    SHW = (CW - 2 * GAP) / 3
    SHH = Inches(0.55)
    SHY = UY + UH + Inches(0.65)
    for i, shop in enumerate(shops):
        sx = ML + i * (SHW + GAP)
        add_rect(s, sx, SHY, SHW, SHH, PANEL, BLUE, 1)
        add_text(s, shop, sx + Inches(0.15), SHY + Inches(0.1),
                 SHW - Inches(0.2), SHH - Inches(0.1),
                 font_size=13, color=TEXT)


def slide_09_bulk(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Moduł 3 — Wyszukiwanie zbiorcze")
    divider(s)

    add_text(s,
             "Idealne gdy masz listę dziesiątek lub setek materiałów do skatalogowania — "
             "np. przy inwentaryzacji lub tworzeniu nowej dokumentacji technicznej.",
             ML, MT, CW, Inches(0.75),
             font_size=15, color=TEXT)

    # Diagram Excel → System → Excel
    GAP  = Inches(0.25)
    FBW  = (CW - 2 * GAP) / 3
    FBH  = Inches(1.9)
    FBY  = MT + Inches(0.9)
    items = [
        ("📊 Plik wejściowy\n(.xlsx)\n\nKolumna:\nopis_materialu", PANEL,  CYAN),
        ("⚙️  IndeksyGSR\n\nPrzetwarza\nwszystkie opisy\nautomatycznie", PANEL2, BLUE),
        ("📊 Plik wynikowy\n(.xlsx)\n\nIndeksy\n+ % pewności", PANEL,  GREEN),
    ]
    xs = [ML + i * (FBW + GAP) for i in range(3)]
    for x, (text, fill, col) in zip(xs, items):
        flow_box(s, x, FBY, FBW, FBH, text, fill, col, 13, col)

    mid = FBY + FBH / 2
    for i in range(2):
        arrow(s, xs[i] + FBW, mid, xs[i + 1], mid)

    # Tabela przykład
    TBY = FBY + FBH + Inches(0.2)
    add_text(s, "Fragment pliku wynikowego:",
             ML, TBY, CW, Inches(0.35),
             font_size=12, bold=True, color=MUTED)

    HDR_H = Inches(0.33)
    TBY2  = TBY + Inches(0.38)
    add_rect(s, ML, TBY2, CW, HDR_H, PANEL2)
    cols_hdr = [
        ("Opis materiału",  0.0,  4.5),
        ("Indeks",          4.6,  2.2),
        ("Nazwa materiału", 6.9,  4.0),
        ("Pewność",        10.9,  1.2),
    ]
    for lbl, xo, ww in cols_hdr:
        add_text(s, lbl, ML + Inches(xo), TBY2 + Inches(0.04),
                 Inches(ww), HDR_H,
                 font_size=10, bold=True, color=MUTED)

    ex_rows = [
        ("śruba M20 ocynkowana ogniowo 80mm", "ŚRUBA-M20-DIN931-OGNIOWO", "Śruba M20×80 DIN 931 ocynk.", "94%", GREEN),
        ("rura stalowa DN50 bezszwowa St35",   "RURA-DN50-ST35-BW",        "Rura stalowa DN50 bezszwowa",  "88%", GREEN),
    ]
    avail = BOT - TBY2 - HDR_H
    RH = avail / len(ex_rows)
    for i, (opis, indeks, nazwa, pct, color) in enumerate(ex_rows):
        ry = TBY2 + HDR_H + i * RH
        add_rect(s, ML, ry, CW, RH - Inches(0.03), PANEL)
        add_text(s, opis,   ML + Inches(0.1),  ry + Inches(0.05), Inches(4.4), RH, font_size=11, color=TEXT)
        add_text(s, indeks, ML + Inches(4.7),  ry + Inches(0.05), Inches(2.1), RH, font_size=11, bold=True, color=CYAN)
        add_text(s, nazwa,  ML + Inches(7.0),  ry + Inches(0.05), Inches(3.8), RH, font_size=11, color=TEXT)
        add_text(s, pct,    ML + CW - Inches(1.1), ry + Inches(0.05), Inches(1.0), RH, font_size=12, bold=True, color=color, align=PP_ALIGN.RIGHT)


def slide_10_stats(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Wyniki projektu")
    divider(s)

    # 4 duże statystyki
    GAP  = Inches(0.2)
    SW4  = (CW - 3 * GAP) / 4
    SH4  = Inches(1.75)
    SY   = MT
    stats = [
        ("68 991",  "indeksów w bazie",         CYAN),
        ("< 1 s",   "czas odpowiedzi",           GREEN),
        ("85%",     "skuteczność scrapowania",   BLUE),
        ("4",       "moduły aplikacji",          YELLOW),
    ]
    for i, (num, lbl, color) in enumerate(stats):
        stat_card(s, ML + i * (SW4 + GAP), SY, SW4, SH4, num, lbl, color)

    # 4 dodatkowe cechy
    IY   = SY + SH4 + Inches(0.2)
    IH   = BOT - IY
    IGAP = Inches(0.2)
    IW   = (CW - IGAP) / 2
    items = [
        ("🌐", "Działa w przeglądarce",    "Nie wymaga instalacji — dostęp z każdego komputera w sieci"),
        ("🔄", "Samouczący się system",    "Każde wyszukiwanie jest rejestrowane i poprawia przyszłe wyniki"),
        ("📈", "Skalowalny",               "Baza może zostać rozszerzona o kolejne katalogi bez przebudowy"),
        ("🔗", "Gotowy do integracji",     "Możliwość podłączenia do systemów ERP / SAP przez interfejs API"),
    ]
    add_rect(s, ML, IY, CW, IH, PANEL, CYAN, 1)
    ITEM_H = IH / 2
    for i, (icon, title, desc) in enumerate(items):
        x = ML + Inches(0.25) + (i % 2) * (IW + IGAP)
        y = IY + (i // 2) * ITEM_H + Inches(0.1)
        add_text(s, f"{icon}  {title}", x, y, IW - Inches(0.2), Inches(0.38),
                 font_size=14, bold=True, color=CYAN)
        add_text(s, desc, x, y + Inches(0.38), IW - Inches(0.2), ITEM_H - Inches(0.55),
                 font_size=12, color=MUTED)


def slide_11_summary(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Podsumowanie")
    divider(s)

    GAP  = Inches(0.25)
    COL  = (CW - GAP) / 2
    COLH = BOT - MT - Inches(0.85)  # zostaw miejsce na stopkę

    # Kolumna lewa
    add_text(s, "Co zostało zrealizowane:",
             ML, MT, COL, Inches(0.42),
             font_size=17, bold=True, color=GREEN)
    done = [
        "✅  System wyszukiwania semantycznego",
        "✅  Baza 68 991 indeksów materiałowych",
        "✅  4 moduły: tekst, URL, zbiorcze, propozycje",
        "✅  Interfejs webowy (przeglądarka)",
        "✅  Automatyczne scrapowanie stron sklepów",
        "✅  Eksport wyników do pliku Excel",
        "✅  Ocena pewności każdego dopasowania",
    ]
    lh = COLH / len(done)
    for i, item in enumerate(done):
        add_text(s, item, ML, MT + Inches(0.5) + i * lh, COL, lh,
                 font_size=14, color=TEXT)

    # Kolumna prawa
    x2 = ML + COL + GAP
    add_text(s, "Możliwości rozwoju:",
             x2, MT, COL, Inches(0.42),
             font_size=17, bold=True, color=BLUE)
    next_steps = [
        "🔷  Integracja z systemami ERP / SAP",
        "🔷  Rozszerzenie na inne katalogi branżowe",
        "🔷  Automatyczne przetwarzanie zamówień",
        "🔷  Aplikacja mobilna",
        "🔷  Moduł zatwierdzania dla działów zakupów",
    ]
    lh2 = COLH / len(next_steps)
    for i, item in enumerate(next_steps):
        add_text(s, item, x2, MT + Inches(0.5) + i * lh2, COL, lh2,
                 font_size=14, color=TEXT)

    # Stopka
    add_rect(s, ML, BOT - Inches(0.65), CW, Inches(0.65), PANEL2)
    add_text(s, "IndeksyGSR  ·  Politechnika Gdańska  ·  2025 / 2026",
             ML, BOT - Inches(0.58), CW, Inches(0.5),
             font_size=13, color=MUTED, align=PP_ALIGN.CENTER)


# ── MAIN ──────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slides = [
        (slide_01_title,    "Tytuł"),
        (slide_02_context,  "Kontekst — indeksy materiałów"),
        (slide_03_problem,  "Problem"),
        (slide_04_goal,     "Cel projektu"),
        (slide_05_how,      "Jak to działa"),
        (slide_06_modules,  "Moduły aplikacji"),
        (slide_07_search_text, "Wyszukiwanie po tekście"),
        (slide_08_url,      "Wyszukiwanie po URL"),
        (slide_09_bulk,     "Wyszukiwanie zbiorcze"),
        (slide_10_stats,    "Wyniki i statystyki"),
        (slide_11_summary,  "Podsumowanie"),
    ]

    print("Generowanie slajdów...")
    for i, (fn, label) in enumerate(slides, 1):
        fn(prs)
        print(f"  {i:2}/{len(slides)}  {label}")

    out = "IndeksyGSR_prezentacja.pptx"
    prs.save(out)
    print(f"\n✅  Zapisano: {out}")


if __name__ == "__main__":
    main()
